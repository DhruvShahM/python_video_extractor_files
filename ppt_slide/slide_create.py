from pptx import Presentation
import tkinter as tk
from tkinter import filedialog, simpledialog

class PPTEditorApp:
    def __init__(self, root):
        self.root = root
        self.root.withdraw()  # Hide main Tkinter window
        self.ppt = None
        self.slides = []
        self.current_slide_index = 0
        self.file_path = filedialog.askopenfilename(
            title="Select a PowerPoint Template",
            filetypes=[("PowerPoint Files", "*.pptx")]
        )
        if not self.file_path:
            exit()
        
        self.ppt = Presentation(self.file_path)
        self.slides = self.ppt.slides
        
        # Ask user how many slides they want to create
        num_new_slides = simpledialog.askinteger("Slide Count", "How many new slides do you want to create?", minvalue=1)
        for _ in range(num_new_slides):
            slide_layout = self.ppt.slide_layouts[1]  # Title and Content layout
            self.ppt.slides.add_slide(slide_layout)
        
        self.slides = self.ppt.slides
        self.num_slides = len(self.slides)
        
        # Set title and subtitle for the first slide
        first_slide = self.slides[0]
        if first_slide.shapes.title:
            title = simpledialog.askstring("Slide Title", "Enter title for the first slide:")
            first_slide.shapes.title.text = title if title else "Default Title"
        
        # Fetch the correct placeholder for subtitle
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape != first_slide.shapes.title:
                subtitle = simpledialog.askstring("Slide Subtitle", "Enter subtitle for the first slide:")
                shape.text = subtitle if subtitle else "Default Subtitle"
                break
        
        self.show_slide_editor()
    
    def show_slide_editor(self):
        self.editor_window = tk.Toplevel(self.root)
        self.editor_window.title(f"Editing Slide {self.current_slide_index + 1}")
        self.editor_window.geometry("600x400")
        
        slide = self.slides[self.current_slide_index]
        
        tk.Label(self.editor_window, text=f"Slide {self.current_slide_index + 1} Title:").pack()
        self.title_var = tk.StringVar()
        title_entry = tk.Entry(self.editor_window, textvariable=self.title_var, width=50)
        title_entry.pack()
        
        tk.Label(self.editor_window, text=f"Slide {self.current_slide_index + 1} Content:").pack()
        self.content_text = tk.Text(self.editor_window, width=50, height=15)
        self.content_text.pack(expand=True, fill='both')
        
        # Load existing slide title
        if slide.shapes.title:
            self.title_var.set(slide.shapes.title.text)
        
        # Find the correct text placeholder
        content_placeholder = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content_placeholder = shape
                break
        
        # Load content if available
        if content_placeholder:
            self.content_text.insert("1.0", content_placeholder.text)
        
        def save_slide():
            """Saves the current slide’s content before switching."""
            if slide.shapes.title:
                slide.shapes.title.text = self.title_var.get()
            if content_placeholder:
                content_placeholder.text = self.content_text.get("1.0", "end-1c")
        
        def next_slide():
            save_slide()
            if self.current_slide_index < self.num_slides - 1:
                self.current_slide_index += 1
                self.editor_window.destroy()
                self.show_slide_editor()
        
        def prev_slide():
            save_slide()
            if self.current_slide_index > 0:
                self.current_slide_index -= 1
                self.editor_window.destroy()
                self.show_slide_editor()
        
        btn_frame = tk.Frame(self.editor_window)
        btn_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=10)
        
        tk.Button(btn_frame, text="Back", command=prev_slide).pack(side=tk.LEFT, padx=5)
        tk.Button(btn_frame, text="Save & Exit", command=self.save_and_exit).pack(side=tk.LEFT, padx=5)
        tk.Button(btn_frame, text="Next", command=next_slide).pack(side=tk.LEFT, padx=5)
    
    def save_and_exit(self):
        output_file = filedialog.asksaveasfilename(
            title="Save Updated Presentation",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")]
        )
        if output_file:
            self.ppt.save(output_file)
        self.editor_window.destroy()
        self.root.quit()

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTEditorApp(root)
    root.mainloop()
