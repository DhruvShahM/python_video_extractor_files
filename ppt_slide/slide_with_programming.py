import os
import tkinter as tk
from tkinter import filedialog, simpledialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup

# Function to select PowerPoint template
def choose_ppt_template():
    root = tk.Tk()
    root.withdraw()
    ppt_path = filedialog.askopenfilename(title="Select a PowerPoint Template", filetypes=[("PowerPoint Files", "*.pptx")])
    return ppt_path

# Function to get programming language from user
def get_language_input():
    root = tk.Tk()
    root.withdraw()
    language = simpledialog.askstring("Programming Language", "Enter the programming language (e.g., python, java, c++):", parent=root)
    return language.lower() if language else "java"  # Default to Java

# Function to get the number of slides from the user
def get_slide_count():
    root = tk.Tk()
    root.withdraw()
    slide_count = simpledialog.askinteger("Number of Slides", "How many slides do you want for the code snippet?", minvalue=1, parent=root)
    return slide_count if slide_count else 1  # Default to 1 if not provided

# Function to get multi-line user input for each slide
def get_code_input(slide_count):
    root = tk.Tk()
    root.withdraw()
    code_snippets = []
    for i in range(slide_count):
        code_text = simpledialog.askstring("Code Input", f"Enter your code snippet for slide {i+1}:", parent=root)
        code_snippets.append(code_text if code_text else "")
    return code_snippets

# Function to get output folder and filename
def get_output_path():
    root = tk.Tk()
    root.withdraw()
    folder_selected = filedialog.askdirectory(title="Select Folder to Save PowerPoint")
    if not folder_selected:
        return None, None
    file_name = simpledialog.askstring("Save As", "Enter output PowerPoint filename (without extension):", parent=root)
    if not file_name:
        return None, None
    return folder_selected, file_name + ".pptx"

# Function to format code using Pygments
def format_code_for_ppt(code, language):
    try:
        lexer = lexers.get_lexer_by_name(language)
    except pygments.util.ClassNotFound:
        lexer = lexers.get_lexer_by_name("java")  # Default to Java if unknown
    formatter = formatters.HtmlFormatter()
    highlighted_code = pygments.highlight(code, lexer, formatter)
    
    # Remove HTML tags
    soup = BeautifulSoup(highlighted_code, "html.parser")
    return soup.get_text()

# Get the PowerPoint template from user
ppt_template = choose_ppt_template()
if not ppt_template or not os.path.exists(ppt_template):
    print("No valid PowerPoint template selected.")
    exit()

# Get programming language from user
language = get_language_input()

# Get the number of slides from the user
slide_count = get_slide_count()

# Get code input for each slide from the user
code_snippets = get_code_input(slide_count)

# Get output folder and filename
output_folder, output_filename = get_output_path()
if not output_folder or not output_filename:
    print("Invalid output path or filename.")
    exit()

# Load the selected PowerPoint template
prs = Presentation(ppt_template)

# Add slides and insert code snippets
for i in range(slide_count):
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Beautify the code
    formatted_code = format_code_for_ppt(code_snippets[i], language)
    
    # Add a textbox for code content
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(5))
    text_frame = textbox.text_frame
    text_frame.text = ""
    text_frame.word_wrap = False  # Preserve indentation and structure

    # Format text inside the textbox
    p = text_frame.add_paragraph()
    p.text = formatted_code
    p.font.name = "Courier New"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Set background color to black for better visibility
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(0, 0, 0)

# Save modified PowerPoint
output_pptx = os.path.join(output_folder, output_filename)
prs.save(output_pptx)
print(f"PPTX saved as {output_pptx}")