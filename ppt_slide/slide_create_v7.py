import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import uuid

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")

        self.prs = Presentation()
        self.slides = []  # Stores (id, title, content, slide_type, font_size, text_color, bg_color)
        self.current_slide = -1
        self.slide_number = 1

        tk.Label(root, text="स्लाइड टाइटल:", font=("Arial", 12)).grid(row=0, column=0, padx=10, pady=5)
        self.title_entry = tk.Entry(root, width=40, font=("Arial", 12))
        self.title_entry.grid(row=0, column=1, padx=10, pady=5)

        self.slide_number_label = tk.Label(root, text=f"स्लाइड नंबर: {self.slide_number}", font=("Arial", 12), fg="blue")
        self.slide_number_label.grid(row=0, column=2, padx=10, pady=5)

        tk.Label(root, text="स्लाइड प्रकार:", font=("Arial", 12)).grid(row=1, column=0, padx=10, pady=5)
        self.slide_type = tk.StringVar(root)
        self.slide_type.set("text")
        tk.OptionMenu(root, self.slide_type, "text", "code", "image + text").grid(row=1, column=1, padx=10, pady=5)

        tk.Label(root, text="अपना कंटेंट डालें:", font=("Arial", 12)).grid(row=2, column=0, padx=10, pady=5)
        self.content_text = tk.Text(root, width=60, height=10, font=("Arial", 12))
        self.content_text.grid(row=2, column=1, padx=10, pady=5)

        tk.Label(root, text="कोड फ़ॉन्ट आकार (pt):", font=("Arial", 12)).grid(row=3, column=0, padx=10, pady=5)
        self.font_size_entry = tk.Entry(root, width=10, font=("Arial", 12))
        self.font_size_entry.grid(row=3, column=1, padx=10, pady=5)
        self.font_size_entry.insert(0, "16")  # Default font size is 16pt

        tk.Label(root, text="कोड टेक्स्ट रंग (#hex):", font=("Arial", 12)).grid(row=4, column=0, padx=10, pady=5)
        self.text_color_entry = tk.Entry(root, width=20, font=("Arial", 12))
        self.text_color_entry.grid(row=4, column=1, padx=10, pady=5)
        self.text_color_entry.insert(0, "#FFFFFF")  # Default text color is white

        tk.Label(root, text="कोड बैकग्राउंड रंग (#hex):", font=("Arial", 12)).grid(row=5, column=0, padx=10, pady=5)
        self.bg_color_entry = tk.Entry(root, width=20, font=("Arial", 12))
        self.bg_color_entry.grid(row=5, column=1, padx=10, pady=5)
        self.bg_color_entry.insert(0, "#000000")  # Default background color is black

        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.grid(row=6, column=1, pady=5)

        tk.Button(root, text="➕ स्लाइड जोड़ें", font=("Arial", 12), command=self.add_slide).grid(row=7, column=0, pady=10)
        tk.Button(root, text="💾 PPT सेव करें", font=("Arial", 12), command=self.save_ppt).grid(row=7, column=1, pady=10)
        tk.Button(root, text="⬅️ पिछली स्लाइड", font=("Arial", 12), command=self.prev_slide).grid(row=8, column=0, pady=10)
        tk.Button(root, text="➡️ अगली स्लाइड", font=("Arial", 12), command=self.next_slide).grid(row=8, column=1, pady=10)
        tk.Button(root, text="🧹 साफ करें", font=("Arial", 12), command=self.clear_fields).grid(row=8, column=2, pady=10)

        tk.Button(root, text="🎨 Image + Text स्लाइड जोड़ें", font=("Arial", 12), command=self.add_image_and_text_slide_from_input).grid(row=9, column=0, pady=10)

    def add_slide(self):
        slide_id = str(uuid.uuid4())  # Generate a unique ID for the slide
        title = self.title_entry.get().strip()
        content = self.content_text.get("1.0", tk.END).strip()
        slide_type = self.slide_type.get()

        font_size = self.font_size_entry.get().strip()
        text_color = self.text_color_entry.get().strip()
        bg_color = self.bg_color_entry.get().strip()

        if not title or not content:
            self.status_label.config(text="⚠️ कृपया टाइटल और कंटेंट भरें!", fg="red")
            return

        try:
            font_size = int(font_size) if font_size else 16  # Default font size 16
        except ValueError:
            font_size = 16

        if not text_color.startswith("#"):
            text_color = "#FFFFFF"  # Default white text color
        if not bg_color.startswith("#"):
            bg_color = "#000000"  # Default black background color

        self.slides.append((slide_id, title, content, slide_type, font_size, text_color, bg_color))
        self.current_slide = len(self.slides) - 1
        self.slide_number += 1
        self.slide_number_label.config(text=f"स्लाइड नंबर: {self.slide_number}")

        if slide_type == "code":
            self.add_code_slide(title, content, font_size, text_color, bg_color)
        elif slide_type == "image + text":
            self.add_image_and_text_slide_from_input()
        else:
            self.add_text_slide(title, content)

        self.status_label.config(text="✅ स्लाइड सफलतापूर्वक जोड़ी गई!", fg="green")
        self.clear_fields()

    def add_text_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

    def add_code_slide(self, title, code, font_size, text_color, bg_color):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        formatted_code = self.format_code_for_ppt(code, "python")
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False

        p = text_frame.add_paragraph()
        p.text = formatted_code
        p.font.name = "Courier New"
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))

    def add_image_and_text_slide_from_input(self):
        title = self.title_entry.get().strip()
        
        # Open file dialog to let the user choose an image file
        image_path = filedialog.askopenfilename(filetypes=[("Image Files", "*.png;*.jpg;*.jpeg")], title="Choose an Image")
        
        if not image_path:
            self.status_label.config(text="⚠️ कृपया एक छवि चुनें!", fg="red")
            return

        text_content = self.content_text.get("1.0", tk.END).strip()

        if not title or not text_content:
            self.status_label.config(text="⚠️ कृपया सभी फ़ील्ड भरें!", fg="red")
            return

        font_size = self.font_size_entry.get().strip()
        text_color = self.text_color_entry.get().strip()
        bg_color = self.bg_color_entry.get().strip()

        try:
            font_size = int(font_size) if font_size else 16
        except ValueError:
            font_size = 16

        if not text_color.startswith("#"):
            text_color = "#FFFFFF"
        if not bg_color.startswith("#"):
            bg_color = "#000000"

        # Call the method to add the slide with the selected image
        self.add_image_and_text_slide(title, image_path, text_content, font_size, text_color, bg_color)

        self.status_label.config(text="✅ स्लाइड सफलतापूर्वक जोड़ी गई!", fg="green")
        self.clear_fields()

    def add_image_and_text_slide(self, title, image_path, text_content, font_size, text_color, bg_color):
        # Choose a slide layout, such as a blank layout or any predefined layout with two columns
        slide_layout = self.prs.slide_layouts[5]  # Choose a blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title

        # Add Image to the slide (First column)
        image_left = Inches(0.5)  # Left position of the image
        image_top = Inches(1.5)   # Top position of the image
        image_width = Inches(4)   # Width of the image
        image_height = Inches(3)  # Height of the image
        slide.shapes.add_picture(image_path, image_left, image_top, image_width, image_height)

        # Add Text to the second column (Second column)
        text_left = Inches(5)    # Left position of the second column (start after the image)
        text_top = Inches(1.5)   # Top position of the text
        text_width = Inches(4)   # Width of the text box
        text_height = Inches(3)  # Height of the text box

        # Create a textbox for the text content
        textbox = slide.shapes.add_textbox(image_left + Inches(4.5), text_top, text_width, text_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        p = text_frame.add_paragraph()
        p.text = text_content
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

        # Set the background color of the textbox
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))

    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def clear_fields(self):
        self.title_entry.delete(0, tk.END)
        self.content_text.delete("1.0", tk.END)
        self.font_size_entry.delete(0, tk.END)
        self.font_size_entry.insert(0, "16")
        self.text_color_entry.delete(0, tk.END)
        self.text_color_entry.insert(0, "#FFFFFF")
        self.bg_color_entry.delete(0, tk.END)
        self.bg_color_entry.insert(0, "#000000")
        self.status_label.config(text="")

    def prev_slide(self):
        if self.current_slide > 0:
            self.current_slide -= 1
            _, title, content, slide_type, font_size, text_color, bg_color = self.slides[self.current_slide]
            self.title_entry.delete(0, tk.END)
            self.title_entry.insert(0, title)
            self.content_text.delete("1.0", tk.END)
            self.content_text.insert(tk.END, content)
            self.font_size_entry.delete(0, tk.END)
            self.font_size_entry.insert(0, font_size)
            self.text_color_entry.delete(0, tk.END)
            self.text_color_entry.insert(0, text_color)
            self.bg_color_entry.delete(0, tk.END)
            self.bg_color_entry.insert(0, bg_color)
            self.slide_number -= 1
            self.slide_number_label.config(text=f"स्लाइड नंबर: {self.slide_number}")

    def next_slide(self):
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            _, title, content, slide_type, font_size, text_color, bg_color = self.slides[self.current_slide]
            self.title_entry.delete(0, tk.END)
            self.title_entry.insert(0, title)
            self.content_text.delete("1.0", tk.END)
            self.content_text.insert(tk.END, content)
            self.font_size_entry.delete(0, tk.END)
            self.font_size_entry.insert(0, font_size)
            self.text_color_entry.delete(0, tk.END)
            self.text_color_entry.insert(0, text_color)
            self.bg_color_entry.delete(0, tk.END)
            self.bg_color_entry.insert(0, bg_color)
            self.slide_number += 1
            self.slide_number_label.config(text=f"स्लाइड नंबर: {self.slide_number}")
            
    def format_code_for_ppt(self, code, language):
        try:
            lexer = lexers.get_lexer_by_name(language)
        except pygments.util.ClassNotFound:
            lexer = lexers.get_lexer_by_name("python")
        formatter = formatters.HtmlFormatter()
        highlighted_code = pygments.highlight(code, lexer, formatter)
        soup = BeautifulSoup(highlighted_code, "html.parser")
        return soup.get_text()            

    def save_ppt(self):
        file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])
        if file_path:
            self.prs.save(file_path)
            self.status_label.config(text=f"✅ PPT सफलतापूर्वक {file_path} में सेव किया गया!", fg="green")

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()
