"""
PowerPoint Generator Pro – v2
────────────────────────────────────────────
Adds:
  •  Separate-file / single-file output modes
  •  Theme Manager (save & load custom themes)
  •  Auto-summary slide toggle
  •  Last-folder memory
Dependencies: python-pptx, markdown2, beautifulsoup4, pygments
Tested: Python 3.9 +
"""
from __future__ import annotations
import os, json, sys, logging, tkinter as tk
from typing import List, Dict, Any
from tkinter import ttk, filedialog, messagebox, colorchooser, font

# ───── third-party ───── #
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import markdown2, pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup

# ───── basic logging ──── #
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-8s | %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
log = logging.getLogger("ppt-gui")

# ───── helpers / constants ───── #
BUILTIN_THEMES: Dict[str, Dict[str, str]] = {
    "Light": {"bg_color": "#FFFFFF", "text_color": "#000000", "font": "Arial"},
    "Dark":  {"bg_color": "#2E2E2E", "text_color": "#FFFFFF", "font": "Arial"},
    "Blue":  {"bg_color": "#1E3A5F", "text_color": "#FFFFFF", "font": "Arial"},
}
CUSTOM_THEME_FILE = "custom_themes.json"


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


# ═══════════════════════════ GUI APP ═══════════════════════════════ #
class PPTApp(tk.Tk):
    def __init__(self) -> None:
        super().__init__()
        self.title("PowerPoint Generator Pro")
        self.geometry("1040x650")
        self.minsize(960, 580)

        # state
        self.files: List[str] = []
        self.prs: Presentation | None = None
        self.slide_titles: List[str] = []
        self.theme: Dict[str, str] = BUILTIN_THEMES["Light"].copy()
        self.custom_themes: Dict[str, Dict[str, str]] = self._load_custom_themes()
        self.last_dir = os.getcwd()

        # tk-vars
        self.v_theme = tk.StringVar(value="Light")
        self.v_bg = tk.StringVar(value=self.theme["bg_color"])
        self.v_fg = tk.StringVar(value=self.theme["text_color"])
        self.v_font = tk.StringVar(value=self.theme["font"])
        self.v_status = tk.StringVar(value="Idle")
        self.v_separate = tk.BooleanVar(value=False)
        self.v_summary = tk.BooleanVar(value=True)

        # build UI
        self._make_menu()
        self._make_toolbar()
        self._make_body()
        self._make_status()

        # shortcuts
        self.bind("<Control-o>", lambda *_: self.add_files())
        self.bind("<Control-s>", lambda *_: self.save_combined())
        self.theme_combo.bind("<<ComboboxSelected>>", self._apply_preset)

    # ───────── UI builders ───────── #
    def _make_menu(self) -> None:
        m = tk.Menu(self)
        fm = tk.Menu(m, tearoff=False)
        fm.add_command(label="Add JSON…      Ctrl+O", command=self.add_files)
        fm.add_command(label="Save combined… Ctrl+S", command=self.save_combined)
        fm.add_separator()
        fm.add_command(label="Exit", command=self.quit)
        m.add_cascade(label="File", menu=fm)

        em = tk.Menu(m, tearoff=False)
        em.add_command(label="Save current theme…", command=self.save_theme_dialog)
        m.add_cascade(label="Edit", menu=em)

        hm = tk.Menu(m, tearoff=False)
        hm.add_command(
            label="About",
            command=lambda: messagebox.showinfo(
                "About", "PowerPoint Generator Pro\n© 2025 Your Name"
            ),
        )
        m.add_cascade(label="Help", menu=hm)
        self.config(menu=m)

    def _make_toolbar(self) -> None:
        tb = ttk.Frame(self, padding=4)
        tb.pack(fill=tk.X)
        ttk.Button(tb, text="➕ Add JSON", command=self.add_files).pack(side=tk.LEFT)
        ttk.Button(tb, text="🚀 Generate", command=self.generate).pack(
            side=tk.LEFT, padx=3
        )
        ttk.Button(tb, text="💾 Save combined", command=self.save_combined).pack(
            side=tk.LEFT, padx=3
        )
        ttk.Separator(tb, orient=tk.VERTICAL).pack(side=tk.LEFT, fill=tk.Y, padx=6)
        ttk.Checkbutton(tb, text="Separate PPTX per JSON", variable=self.v_separate).pack(
            side=tk.LEFT
        )
        ttk.Checkbutton(tb, text="Add summary slide", variable=self.v_summary).pack(
            side=tk.LEFT, padx=4
        )

    def _make_body(self) -> None:
        body = ttk.Frame(self)
        body.pack(fill=tk.BOTH, expand=True, padx=6, pady=6)
        body.columnconfigure(1, weight=1)
        body.rowconfigure(0, weight=1)

        # left – queue & theme
        left = ttk.Frame(body)
        left.grid(row=0, column=0, sticky="nsw")
        # queue
        ttk.Label(left, text="File queue").pack(anchor="w")
        self.lb_files = tk.Listbox(left, width=35, height=14)
        self.lb_files.pack()
        row = ttk.Frame(left)
        row.pack(pady=3)
        ttk.Button(row, text="⬆", width=3, command=lambda: self._move(-1)).pack(
            side=tk.LEFT
        )
        ttk.Button(row, text="⬇", width=3, command=lambda: self._move(1)).pack(
            side=tk.LEFT, padx=2
        )
        ttk.Button(row, text="🗑", width=3, command=self._remove).pack(side=tk.LEFT)
        ttk.Separator(left, orient=tk.HORIZONTAL).pack(fill=tk.X, pady=6)

        # theme area
        ttk.Label(left, text="Theme").pack(anchor="w")
        all_theme_names = list(BUILTIN_THEMES) + list(self.custom_themes)
        self.theme_combo = ttk.Combobox(
            left, textvariable=self.v_theme, values=all_theme_names, state="readonly"
        )
        self.theme_combo.pack(fill=tk.X)
        # colour + font pickers
        self._colour_row(left, "Background", self.v_bg)
        self._colour_row(left, "Text colour", self.v_fg)
        ttk.Label(left, text="Font").pack(anchor="w", pady=(6, 0))
        ttk.Combobox(
            left, textvariable=self.v_font, values=font.families(), state="readonly"
        ).pack(fill=tk.X, pady=(0, 6))
        ttk.Button(left, text="💾 Save theme", command=self.save_theme_dialog).pack(
            fill=tk.X
        )

        # right – preview
        right = ttk.Frame(body)
        right.grid(row=0, column=1, sticky="nsew", padx=(10, 0))
        right.columnconfigure(0, weight=1)
        right.rowconfigure(1, weight=1)
        ttk.Label(right, text="Slide titles preview").grid(row=0, column=0, sticky="w")
        preview_frame = ttk.Frame(right)
        preview_frame.grid(row=1, column=0, sticky="nsew")
        preview_frame.columnconfigure(0, weight=1)
        preview_frame.rowconfigure(0, weight=1)
        self.lb_preview = tk.Listbox(preview_frame)
        self.lb_preview.grid(row=0, column=0, sticky="nsew")
        sb = ttk.Scrollbar(
            preview_frame, orient="vertical", command=self.lb_preview.yview
        )
        sb.grid(row=0, column=1, sticky="ns")
        self.lb_preview.configure(yscrollcommand=sb.set)

    def _make_status(self) -> None:
        st = ttk.Frame(self, relief=tk.SUNKEN)
        st.pack(fill=tk.X, side=tk.BOTTOM)
        ttk.Label(st, textvariable=self.v_status).pack(side=tk.LEFT, padx=6)
        self.progress = ttk.Progressbar(st, maximum=100)
        self.progress.pack(side=tk.RIGHT, fill=tk.X, expand=True, padx=6, pady=4)

    # ───────── helpers ───────── #
    def _colour_row(self, parent, label, var):
        ttk.Label(parent, text=label).pack(anchor="w", pady=(6, 0))
        fr = ttk.Frame(parent)
        fr.pack(fill=tk.X)
        ttk.Entry(fr, textvariable=var, state="readonly", width=12).pack(side=tk.LEFT)
        ttk.Button(fr, text="🎨", command=lambda: self._pick_colour(var)).pack(
            side=tk.LEFT, padx=3
        )

    def _pick_colour(self, var):
        c = colorchooser.askcolor(color=var.get())[1]
        if c:
            var.set(c)

    def _apply_preset(self, *_):
        name = self.v_theme.get()
        src = self.custom_themes.get(name) or BUILTIN_THEMES.get(name)
        if src:
            self.theme.update(src)
            self.v_bg.set(src["bg_color"])
            self.v_fg.set(src["text_color"])
            self.v_font.set(src["font"])

    def _refresh_queue(self):
        self.lb_files.delete(0, tk.END)
        for p in self.files:
            self.lb_files.insert(tk.END, os.path.basename(p))

    def _move(self, delta: int):
        sel = self.lb_files.curselection()
        if not sel:
            return
        i = sel[0]
        j = i + delta
        if 0 <= j < len(self.files):
            self.files[i], self.files[j] = self.files[j], self.files[i]
            self._refresh_queue()
            self.lb_files.select_set(j)

    def _remove(self):
        sel = self.lb_files.curselection()
        if not sel:
            return
        self.files.pop(sel[0])
        self._refresh_queue()

    def _load_custom_themes(self) -> Dict[str, Dict[str, str]]:
        if not os.path.exists(CUSTOM_THEME_FILE):
            return {}
        try:
            with open(CUSTOM_THEME_FILE, encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            log.error("Failed loading custom themes: %s", e)
            return {}

    def _save_custom_themes(self):
        try:
            with open(CUSTOM_THEME_FILE, "w", encoding="utf-8") as f:
                json.dump(self.custom_themes, f, indent=2)
        except Exception as e:
            log.error("Saving themes failed: %s", e)

    def save_theme_dialog(self):
        name = tk.simpledialog.askstring("Save theme", "Give your theme a name:")
        if not name:
            return
        all_names = list(BUILTIN_THEMES) + list(self.custom_themes)
        if name in all_names:
            messagebox.showerror("Duplicate", "Theme name already exists.")
            return
        self.custom_themes[name] = {
            "bg_color": self.v_bg.get(),
            "text_color": self.v_fg.get(),
            "font": self.v_font.get(),
        }
        self._save_custom_themes()
        self.theme_combo["values"] = list(BUILTIN_THEMES) + list(self.custom_themes)
        messagebox.showinfo("Saved", f"Theme “{name}” saved!")

    # ───────── file actions ───────── #
    def add_files(self):
        paths = filedialog.askopenfilenames(
            title="Select JSON slide files",
            filetypes=[("JSON files", "*.json")],
            initialdir=self.last_dir,
        )
        if not paths:
            return
        self.last_dir = os.path.dirname(paths[0])
        self.files.extend(paths)
        self._refresh_queue()

    # ═════════ generation ═════════ #
    def generate(self):
        if not self.files:
            messagebox.showwarning("No files", "Add at least one JSON file first.")
            return

        # separate vs combined
        if self.v_separate.get():
            out_dir = filedialog.askdirectory(
                mustexist=True,
                title="Choose output folder for individual PPTX files",
                initialdir=self.last_dir,
            )
            if not out_dir:
                return
            self.last_dir = out_dir
            self._generate_separate(out_dir)
        else:
            self._generate_combined()
        self.progress["value"] = 0

    # -- combined deck
    def _generate_combined(self):
        self._reset_preview()
        self.prs = Presentation()
        for ix, path in enumerate(self.files, 1):
            self._process_json(path)
            self.progress["value"] = ix / len(self.files) * 100
            self.update_idletasks()
        if self.v_summary.get():
            self._add_summary_slide(self.slide_titles)
        self.v_status.set("Combined presentation ready – use “Save combined”")

    # -- separate decks
    def _generate_separate(self, out_dir: str):
        for ix, path in enumerate(self.files, 1):
            self.v_status.set(f"Processing {os.path.basename(path)} ({ix}/{len(self.files)})")
            self.update_idletasks()
            prs = Presentation()
            titles_capture: List[str] = []
            self._process_json(path, prs, titles_capture)
            if self.v_summary.get():
                self._add_summary_slide(titles_capture, prs)
            name = os.path.splitext(os.path.basename(path))[0] + ".pptx"
            try:
                prs.save(os.path.join(out_dir, name))
            except Exception as e:
                messagebox.showerror("Save error", f"Could not save {name}\n{e}")
                return
            self.progress["value"] = ix / len(self.files) * 100
            self.update_idletasks()
        messagebox.showinfo("Done", f"Individual files saved to\n{out_dir}")
        self.v_status.set("All individual decks saved")

    # -- JSON ⇒ slides ---------------- #
    def _process_json(
        self,
        path: str,
        prs: Presentation | None = None,
        title_capture: List[str] | None = None,
    ):
        target_prs = prs or self.prs  # combined uses self.prs
        assert target_prs is not None
        try:
            with open(path, encoding="utf-8") as f:
                data: List[Dict[str, Any]] = json.load(f)
        except Exception as e:
            messagebox.showerror("JSON error", f"{path}\n{e}")
            raise

        # first slide of this file becomes a title slide
        self._add_title_slide(target_prs, data[0].get("title", "Presentation"))
        for slide in data:
            t = slide.get("title", "Untitled")
            ct = slide.get("content", "")
            st = slide.get("slide_type", "text")
            nt = slide.get("notes", "")
            self.slide_titles.append(t)
            if title_capture is not None:
                title_capture.append(t)

            if st == "code":
                self._add_code_slide(target_prs, t, ct, nt)
            elif st == "table" and isinstance(ct, list):
                self._add_table_slide(target_prs, t, ct, nt)
            else:
                self._add_text_slide(target_prs, t, ct, nt)

    # ═════════ slide builders ═════════ #
    def _add_title_slide(self, prs: Presentation, title: str):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title
        s.placeholders[1].text = "Auto-generated presentation"

    def _add_text_slide(self, prs: Presentation, title: str, content: str, notes: str):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for line in str(content).split("\n"):
            p = tf.add_paragraph()
            p.text = line
        self._add_notes(s, notes)

    def _add_table_slide(
        self, prs: Presentation, title: str, rows: List[Dict[str, Any]], notes: str
    ):
        if not rows:
            return
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title
        cols = len(rows[0])
        table = s.shapes.add_table(
            len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
        ).table
        # headers
        for c, h in enumerate(rows[0].keys()):
            cell = table.cell(0, c)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
        # data rows
        for r, row in enumerate(rows, 1):
            for c, v in enumerate(row.values()):
                table.cell(r, c).text = str(v)
        self._add_notes(s, notes)

    def _add_code_slide(
        self, prs: Presentation, title: str, code: str, notes: str, max_lines=18
    ):
        lines = code.split("\n")
        for i in range(0, len(lines), max_lines):
            chunk = "\n".join(lines[i : i + max_lines])
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = title if i == 0 else f"{title} (contd.)"
            box = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(9), Inches(5))
            tf = box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = self._highlight(chunk)
            p.font.name = "Courier New"
            p.font.size = Pt(14)
            # colours
            p.font.color.rgb = RGBColor(*hex_to_rgb(self.v_fg.get()))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*hex_to_rgb(self.v_bg.get()))
            self._add_notes(s, notes)

    def _add_summary_slide(
        self, titles: List[str], prs: Presentation | None = None, notes: str = ""
    ):
        target = prs or self.prs
        if not target:
            return
        s = target.slides.add_slide(target.slide_layouts[1])
        s.shapes.title.text = "Summary"
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for t in titles:
            p = tf.add_paragraph()
            p.text = t
            p.level = 0
        self._add_notes(s, notes)

    def _add_notes(self, slide, txt: str):
        if not txt:
            return
        slide.notes_slide.notes_text_frame.text = txt

    def _highlight(self, code: str) -> str:
        try:
            lexer = lexers.get_lexer_by_name("python")
        except pygments.util.ClassNotFound:
            lexer = lexers.get_lexer_by_name("text")
        fmt = formatters.HtmlFormatter()
        soup = BeautifulSoup(pygments.highlight(code, lexer, fmt), "html.parser")
        return soup.get_text()

    # ═════════ save combined ═════════ #
    def save_combined(self):
        if not self.prs:
            messagebox.showwarning(
                "Nothing to save", "Create a combined presentation first."
            )
            return
        path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialdir=self.last_dir,
            title="Save combined presentation as…",
        )
        if not path:
            return
        try:
            self.prs.save(path)
            self.last_dir = os.path.dirname(path)
        except Exception as e:
            messagebox.showerror("Save error", str(e))
            return
        messagebox.showinfo("Saved", f"PPTX written to\n{path}")
        self.v_status.set(f"Saved: {os.path.basename(path)}")

    # ═════════ misc ═════════ #
    def _reset_preview(self):
        self.slide_titles.clear()
        self.lb_preview.delete(0, tk.END)
        self.progress["value"] = 0
        self.v_status.set("Generating…")


# ───── main ───── #
if __name__ == "__main__":
    app = PPTApp()
    app.mainloop()
