import os
import pandas as pd
import ollama
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

# 🟢 Step 1: Load Topics from CSV
topics_file = os.path.join(os.path.dirname(__file__), "topics.csv")

if not os.path.exists(topics_file):
    print("❌ topics.csv file not found! Please create one.")
    exit()

df = pd.read_csv(topics_file)

# 🟢 Step 2: Create a New PowerPoint Presentation
ppt = Presentation()

def generate_text_ollama(prompt):
    """Generate AI-based text using Llama3.2 locally via Ollama."""
    response = ollama.chat(model="llama3.2:latest", messages=[{"role": "user", "content": prompt}])
    return response['message']['content']

# 🟢 Step 3: Generate Slides
for _, row in df.iterrows():
    topic = row["Topic"]

    print(f"🔹 Generating content for: {topic}...")

    # Get AI-generated slide content
    slide_content = generate_text_ollama(f"Write a short summary about {topic} in bullet points.")

    # Add a new slide (Title & Content Layout)
    slide_layout = ppt.slide_layouts[1]  # Title & Content Layout
    slide = ppt.slides.add_slide(slide_layout)

    # Set Slide Title
    slide.shapes.title.text = topic

    # Set Slide Content
    content_placeholder = slide.placeholders[1]  # Text placeholder
    content_placeholder.text = slide_content

# 🟢 Step 4: Save the PowerPoint File
output_dir = "output"
os.makedirs(output_dir, exist_ok=True)

timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
output_path = os.path.join(output_dir, f"presentation_{timestamp}.pptx")

ppt.save(output_path)

print(f"✅ PowerPoint saved at: {output_path}")
