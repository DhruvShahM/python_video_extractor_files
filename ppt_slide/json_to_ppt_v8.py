import os
import json
import tkinter as tk
from tkinter import filedialog, ttk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import textwrap
import markdown2
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

THEMES = {
    "Light": {"bg_color": "#FFFFFF", "text_color": "#000000"},
    "Dark": {"bg_color": "#2E2E2E", "text_color": "#FFFFFF"},
    "Blue": {"bg_color": "#1E3A5F", "text_color": "#FFFFFF"}
}

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")
        self.prs = Presentation()

        self.theme_var = tk.StringVar(value="Light")
        ttk.Label(root, text="Select Theme:", font=("Arial", 12)).pack()
        self.theme_dropdown = ttk.Combobox(root, textvariable=self.theme_var, values=list(THEMES.keys()))
        self.theme_dropdown.pack()

        tk.Button(root, text="📂 Load JSON & Create PPT", font=("Arial", 12), command=self.load_json).pack(pady=20)
        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.pack()

    def add_title_slide(self, title):
        slide_layout = self.prs.slide_layouts[0]  # Title Slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = "Auto-generated Presentation"

    # def add_bullet_slide(self, title, content):
    #     slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
    #     slide = self.prs.slides.add_slide(slide_layout)
    #     title_placeholder = slide.shapes.title
    #     content_placeholder = slide.shapes.placeholders[1]

    #     title_placeholder.text = title

    #     # Check if content is a list
    #     if isinstance(content, list):
    #         for bullet in content:
    #             if isinstance(bullet, dict):  # Convert dict to string safely
    #                 bullet = ", ".join(f"{k}: {v}" for k, v in bullet.items())
                
    #             p = content_placeholder.text_frame.add_paragraph()
    #             p.text = str(bullet)  # Ensure text format
    #     else:
    #         for bullet in str(content).split("\n"):
    #             p = content_placeholder.text_frame.add_paragraph()
    #             p.text = bullet



    def add_summary_slide(self, titles):
        slide_layout = self.prs.slide_layouts[1]  # Title + Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Summary"

        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()

        for title in titles:
            p = text_frame.add_paragraph()
            p.text = title
            p.level = 0  # Bullet point    

    def add_table_slide(self, title, table_data):
        slide_layout = self.prs.slide_layouts[5]  # Title Only Layout
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 0
        
        if rows == 0 or cols == 0:
            return

        table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)).table
        
        # Add headers
        for col_idx, header in enumerate(table_data[0].keys()):
            table.cell(0, col_idx).text = header
            table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
        
        # Add data rows
        for row_idx, row_data in enumerate(table_data):
            for col_idx, (key, value) in enumerate(row_data.items()):
                table.cell(row_idx + 1, col_idx).text = str(value)        
    
    def load_json(self):
        file_paths = filedialog.askopenfilenames(filetypes=[("JSON Files", "*.json")], title="Select JSON Files")
        if not file_paths:
            return

        for file_path in file_paths:
            with open(file_path, "r", encoding="utf-8") as file:
                slides_data = json.load(file)

            if slides_data:
                self.prs = Presentation()  # Create a new presentation for each file
                self.add_title_slide(slides_data[0].get("title", "Presentation"))

                theme = THEMES.get(self.theme_var.get(), THEMES["Light"])
                slide_titles = []

                for slide in slides_data:
                    title = slide.get("title", "Untitled")
                    content = slide.get("content", "")
                    slide_type = slide.get("slide_type", "text")

                    slide_titles.append(title)

                    if slide_type == "code":
                        self.add_code_slide(title, content, theme["text_color"], theme["bg_color"])
                    elif slide_type == "table" and isinstance(content, list):
                        self.add_table_slide(title, content)
                    else:
                        self.add_bullet_slide(title, content)

                self.add_summary_slide(slide_titles)

                save_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                        filetypes=[("PowerPoint Files", "*.pptx")],
                                                        title=f"Save PPT for {os.path.basename(file_path)}")
                if save_path:
                    self.prs.save(save_path)
                    self.status_label.config(text=f"✅ PPT saved: {save_path}", fg="blue")

    def add_text_slide(self, title, content):
        """Convert Markdown content to formatted PowerPoint slides."""
        max_words_per_slide = 300  # Adjust for readability

        # Convert Markdown to HTML
        html_content = markdown2.markdown(content)

        # Parse HTML and extract text
        soup = BeautifulSoup(html_content, "html.parser")

        # Split text into slides for readability
        words = soup.get_text().split()
        
        for i in range(0, len(words), max_words_per_slide):
            slide_layout = self.prs.slide_layouts[1]  # Title + Content layout
            slide = self.prs.slides.add_slide(slide_layout)

            if slide.shapes.title:
                slide.shapes.title.text = title if i == 0 else f"{title} (contd.)"

            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            p = text_frame.add_paragraph()
            p.text = " ".join(words[i:i + max_words_per_slide])
            p.font.name = "Arial"
            p.font.size = Pt(20)  # Adjust size for readability



    def add_code_slide(self, title, code, text_color, bg_color):
                max_lines_per_slide = 20  # Adjust as needed
                lines = code.split("\n")
                
                for i in range(0, len(lines), max_lines_per_slide):
                    slide_layout = self.prs.slide_layouts[5]
                    slide = self.prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = title if i == 0 else f"{title} (contd.)"

                    code_chunk = "\n".join(lines[i:i+max_lines_per_slide])
                    formatted_code = self.format_code_for_ppt(code_chunk, "java")

                    font_size = self.calculate_dynamic_font_size(code_chunk)
                    
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = False

                    p = text_frame.add_paragraph()
                    p.text = formatted_code
                    p.font.name = "Courier New"
                    p.font.size = Pt(font_size)
                    p.alignment = PP_ALIGN.LEFT
                    p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))


                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))

    def format_code_for_ppt(self, code, language):
            try:
                lexer = lexers.get_lexer_by_name(language)
            except pygments.util.ClassNotFound:
                lexer = lexers.get_lexer_by_name("python")
            formatter = formatters.HtmlFormatter()
            highlighted_code = pygments.highlight(code, lexer, formatter)
            soup = BeautifulSoup(highlighted_code, "html.parser")
            return soup.get_text()

    def calculate_dynamic_font_size(self, code):
            line_count = code.count('\n') + 1
            if line_count < 10:
                return 24
            elif line_count < 20:
                return 18
            elif line_count < 30:
                return 14
            else:
                return 12

    def hex_to_rgb(self, hex_color):            
            hex_color = hex_color.lstrip("#")
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def add_bullet_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.shapes.placeholders[1]

        title_placeholder.text = title

        text_frame = content_placeholder.text_frame
        text_frame.clear()

        if isinstance(content, list):
            for bullet in content:
                p = text_frame.add_paragraph()
                self.apply_markdown_formatting(p, bullet)
        else:
            for line in str(content).split("\n"):
                p = text_frame.add_paragraph()
                self.apply_markdown_formatting(p, line)

        text_frame.word_wrap = True

    def apply_markdown_formatting(self, paragraph, text):
        """Processes markdown-like formatting in JSON text."""
        
        # Process **bold**
        if "**" in text:
            bold_sections = text.split("**")
            for i, section in enumerate(bold_sections):
                run = paragraph.add_run()
                run.text = section
                if i % 2 == 1:
                    run.font.bold = True
        else:
            paragraph.text = text

        # Process ### Headings
        if text.startswith("### "):
            paragraph.text = text.replace("### ", "")
            paragraph.font.bold = True
            paragraph.font.size = Pt(24)  # Larger font for headers
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        else:
            paragraph.font.size = Pt(18)  # Normal text size

        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Set text color 

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()