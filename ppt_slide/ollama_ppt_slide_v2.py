import ollama
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Generate code from Ollama
response = ollama.chat(model="codellama", messages=[{"role": "user", "content": "What is a spring boot?"}])
generated_code = response["message"]["content"]

# Generate explanation from Ollama
explanation_response = ollama.chat(model="codellama", messages=[{"role": "user", "content": "Explain this Java function step by step"}])
code_explanation = explanation_response["message"]["content"]

# Create PowerPoint Presentation
ppt = Presentation()

# ------ Slide 1: Code ------ #
slide1 = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide

# Add title
title_shape = slide1.shapes.title
if title_shape:
    title_shape.text = "Generated Java Code"

# Add textbox for code
left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(5)
textbox = slide1.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Format code text
p = text_frame.add_paragraph()
p.text = generated_code
p.font.name = "Courier New"  # Monospace font
p.font.size = Inches(0.3)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White text
textbox.fill.solid()
textbox.fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark background

# ------ Slide 2: Explanation ------ #
slide2 = ppt.slides.add_slide(ppt.slide_layouts[5])  # Blank slide

# Add title
title_shape = slide2.shapes.title
if title_shape:
    title_shape.text = "Code Explanation"

# Add textbox for explanation
textbox = slide2.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Add explanation text
p = text_frame.add_paragraph()
p.text = code_explanation
p.font.name = "Arial"
p.font.size = Inches(0.3)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
ppt.save("Generated_Code_Presentation.pptx")

print("PowerPoint file saved successfully as 'Generated_Code_Presentation.pptx'")
