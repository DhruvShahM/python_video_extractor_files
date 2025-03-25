import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import uuid
import json  # Import json module for handling JSON files
import pandas as pd  # Added pandas for CSV handling

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")

        self.prs = Presentation()
        self.slides = []  # Stores (id, title, content, slide_type, font_color, font_size, background_color)
        self.current_slide = -1
        self.slide_number = 1

        tk.Label(root, text="JSON फ़ाइल चुनें:", font=("Arial", 12)).grid(row=0, column=0, padx=10, pady=5)
        tk.Button(root, text="JSON फ़ाइल चुनें", font=("Arial", 12), command=self.load_json).grid(row=0, column=1, padx=10, pady=5)

        self.slide_number_label = tk.Label(root, text=f"स्लाइड नंबर: {self.slide_number}", font=("Arial", 12), fg="blue")
        self.slide_number_label.grid(row=1, column=0, padx=10, pady=5)

        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.grid(row=2, column=1, pady=5)

        tk.Button(root, text="💾 PPT सेव करें", font=("Arial", 12), command=self.save_ppt).grid(row=3, column=1, pady=10)

    def load_json(self):
        # Open file dialog to select the JSON file
        json_path = filedialog.askopenfilename(filetypes=[("JSON Files", "*.json")], title="JSON फ़ाइल चुनें")
        if json_path:
            self.read_json_and_create_slides(json_path)
            self.status_label.config(text=f"✅ JSON से स्लाइड्स जोड़ी गईं!", fg="green")

    def read_json_and_create_slides(self, json_path):
        try:
            # Read the JSON file
            with open(json_path, "r", encoding="utf-8") as file:
                data = json.load(file)
                for item in data:
                    # Assuming the JSON structure is a list of dictionaries with 'title', 'content', 'slide_type'
                    title = item.get("title")
                    content = item.get("content")
                    slide_type = item.get("slide_type")

                    if not title or not content or not slide_type:
                        self.status_label.config(text="⚠️ Invalid JSON structure.", fg="red")
                        return

                    slide_id = str(uuid.uuid4())  # Generate a unique ID for the slide
                    self.slides.append((slide_id, title, content, slide_type))
                    if slide_type == "code":
                        self.add_code_slide(title, content)
                    else:
                        self.add_text_slide(title, content)

                    self.slide_number += 1
                    self.slide_number_label.config(text=f"स्लाइड नंबर: {self.slide_number}")
        except Exception as e:
            self.status_label.config(text=f"Error reading JSON: {e}", fg="red")
            print(f"Error reading JSON: {e}")

    def add_text_slide(self, title, content, font_color='000000', font_size=18, background_color='FFFFFF'):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        text_placeholder = slide.shapes.placeholders[1]
        text_placeholder.text = content

        self.apply_text_style(text_placeholder, font_color, font_size)
        self.apply_background_color(slide, background_color)

    def add_code_slide(self, title, code, font_color='000000', font_size=18, background_color='FFFFFF'):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        formatted_code = self.format_code_for_ppt(code, "python")
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False

        p = text_frame.add_paragraph()
        p.text = formatted_code
        p.font.name = "Courier New"
        p.font.size = Pt(int(font_size))  # Set font size from CSV
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = self.hex_to_rgb(font_color)

        textbox.fill.solid()
        textbox.fill.fore_color.rgb = self.hex_to_rgb(background_color)

    def apply_text_style(self, text_placeholder, font_color, font_size):
        text_frame = text_placeholder.text_frame
        text_frame.paragraphs[0].font.size = Pt(int(font_size))  # Apply font size
        text_frame.paragraphs[0].font.color.rgb = self.hex_to_rgb(font_color)

    def apply_background_color(self, slide, background_color):
        slide_background = slide.background
        slide_background.fill.solid()
        slide_background.fill.fore_color.rgb = self.hex_to_rgb(background_color)

    def format_code_for_ppt(self, code, language):
        try:
            lexer = lexers.get_lexer_by_name(language)
        except pygments.util.ClassNotFound:
            lexer = lexers.get_lexer_by_name("python")
        formatter = formatters.HtmlFormatter()
        highlighted_code = pygments.highlight(code, lexer, formatter)
        soup = BeautifulSoup(highlighted_code, "html.parser")
        return soup.get_text()

    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip("#")
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def save_ppt(self):
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")], title="PPT को सेव करने के लिए स्थान चुनें")
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(text=f"✅ PPT सेव हो गई: {save_path}", fg="blue")
            self.root.quit()  # एप्लिकेशन बंद करें


if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()
