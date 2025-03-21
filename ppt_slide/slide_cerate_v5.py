import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import uuid

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")

        self.prs = Presentation()
        self.slides = []  # Stores (id, title, content, slide_type)
        self.current_slide = -1
        self.slide_number = 1

        tk.Label(root, text="स्लाइड टाइटल:", font=("Arial", 12)).grid(row=0, column=0, padx=10, pady=5)
        self.title_entry = tk.Entry(root, width=40, font=("Arial", 12))
        self.title_entry.grid(row=0, column=1, padx=10, pady=5)

        self.slide_number_label = tk.Label(root, text=f"स्लाइड नंबर: {self.slide_number}", font=("Arial", 12), fg="blue")
        self.slide_number_label.grid(row=0, column=2, padx=10, pady=5)

        tk.Label(root, text="स्लाइड प्रकार:", font=("Arial", 12)).grid(row=1, column=0, padx=10, pady=5)
        self.slide_type = tk.StringVar(root)
        self.slide_type.set("text")
        tk.OptionMenu(root, self.slide_type, "text", "code").grid(row=1, column=1, padx=10, pady=5)

        tk.Label(root, text="अपना कंटेंट डालें:", font=("Arial", 12)).grid(row=2, column=0, padx=10, pady=5)
        self.content_text = tk.Text(root, width=60, height=10, font=("Arial", 12))
        self.content_text.grid(row=2, column=1, padx=10, pady=5)

        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.grid(row=3, column=1, pady=5)

        tk.Button(root, text="➕ स्लाइड जोड़ें", font=("Arial", 12), command=self.add_slide).grid(row=5, column=0, pady=10)
        tk.Button(root, text="💾 PPT सेव करें", font=("Arial", 12), command=self.save_ppt).grid(row=5, column=1, pady=10)
        tk.Button(root, text="⬅️ पिछली स्लाइड", font=("Arial", 12), command=self.prev_slide).grid(row=6, column=0, pady=10)
        tk.Button(root, text="➡️ अगली स्लाइड", font=("Arial", 12), command=self.next_slide).grid(row=6, column=1, pady=10)
        tk.Button(root, text="🧹 साफ करें", font=("Arial", 12), command=self.clear_fields).grid(row=6, column=2, pady=10)

    def add_slide(self):
        slide_id = str(uuid.uuid4())  # Generate a unique ID for the slide
        title = self.title_entry.get().strip()
        content = self.content_text.get("1.0", tk.END).strip()
        slide_type = self.slide_type.get()

        if not title or not content:
            self.status_label.config(text="⚠️ कृपया टाइटल और कंटेंट भरें!", fg="red")
            return

        self.slides.append((slide_id, title, content, slide_type))
        self.current_slide = len(self.slides) - 1
        self.slide_number += 1
        self.slide_number_label.config(text=f"स्लाइड नंबर: {self.slide_number}")

        if slide_type == "code":
            self.add_code_slide(title, content)
        else:
            self.add_text_slide(title, content)

        self.status_label.config(text="✅ स्लाइड सफलतापूर्वक जोड़ी गई!", fg="green")
        self.clear_fields()

    def add_text_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

    def add_code_slide(self, title, code):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        formatted_code = self.format_code_for_ppt(code, "python")
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False

        p = text_frame.add_paragraph()
        p.text = formatted_code
        p.font.name = "Courier New"
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(255, 255, 255)

        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(0, 0, 0)

    def format_code_for_ppt(self, code, language):
        try:
            lexer = lexers.get_lexer_by_name(language)
        except pygments.util.ClassNotFound:
            lexer = lexers.get_lexer_by_name("python")
        formatter = formatters.HtmlFormatter()
        highlighted_code = pygments.highlight(code, lexer, formatter)
        soup = BeautifulSoup(highlighted_code, "html.parser")
        return soup.get_text()

    def clear_fields(self):
        self.title_entry.delete(0, tk.END)
        self.content_text.delete("1.0", tk.END)
        self.status_label.config(text="")

    def prev_slide(self):
        if self.current_slide > 0:
            self.current_slide -= 1
            _, title, content, slide_type = self.slides[self.current_slide]
            
            self.title_entry.delete(0, tk.END)
            self.title_entry.insert(0, title)
            
            self.content_text.delete("1.0", tk.END)
            self.content_text.insert("1.0", content)
            
            self.slide_type.set(slide_type)  # स्लाइड प्रकार अपडेट करें

    def next_slide(self):
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            _, title, content, slide_type = self.slides[self.current_slide]
            
            self.title_entry.delete(0, tk.END)
            self.title_entry.insert(0, title)
            
            self.content_text.delete("1.0", tk.END)
            self.content_text.insert("1.0", content)
            
            self.slide_type.set(slide_type)  # स्लाइड प्रकार अपडेट करें


    def save_ppt(self):
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")], title="PPT को सेव करने के लिए स्थान चुनें")
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(text=f"✅ PPT सेव हो गई: {save_path}", fg="blue")
            self.root.quit()  # एप्लिकेशन बंद करें


if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()
