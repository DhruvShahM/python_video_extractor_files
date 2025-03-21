import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches
import textwrap

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")

        self.prs = Presentation()
        self.slides = []
        self.current_slide = -1  

        # Title Label
        tk.Label(root, text="स्लाइड टाइटल:", font=("Arial", 12)).grid(row=0, column=0, padx=10, pady=5)
        self.title_entry = tk.Entry(root, width=40, font=("Arial", 12))
        self.title_entry.grid(row=0, column=1, padx=10, pady=5)

        # Type Selection
        tk.Label(root, text="स्लाइड प्रकार:", font=("Arial", 12)).grid(row=1, column=0, padx=10, pady=5)
        self.slide_type = tk.StringVar(root)
        self.slide_type.set("text")
        tk.OptionMenu(root, self.slide_type, "text", "code").grid(row=1, column=1, padx=10, pady=5)

        # Content Box
        tk.Label(root, text="अपना कंटेंट डालें:", font=("Arial", 12)).grid(row=2, column=0, padx=10, pady=5)
        self.content_text = tk.Text(root, width=60, height=10, font=("Arial", 12))
        self.content_text.grid(row=2, column=1, padx=10, pady=5)

        # Status Label
        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.grid(row=3, column=1, pady=5)

        # Navigation Buttons (Previous, Next)
        self.prev_button = tk.Button(root, text="⬅️ पिछला", font=("Arial", 12), command=self.previous_slide, state=tk.DISABLED)
        self.prev_button.grid(row=4, column=0, pady=10)

        self.next_button = tk.Button(root, text="➡️ अगला", font=("Arial", 12), command=self.next_slide, state=tk.DISABLED)
        self.next_button.grid(row=4, column=1, pady=10)

        # Add Slide Button
        tk.Button(root, text="➕ स्लाइड जोड़ें", font=("Arial", 12), command=self.add_slide).grid(row=5, column=0, pady=10)

        # Exit & Save Button
        tk.Button(root, text="💾 PPT सेव करें और बाहर निकलें", font=("Arial", 12), command=self.save_and_exit).grid(row=5, column=1, pady=10)

    def add_slide(self):
        """नया स्लाइड जोड़ें"""
        title = self.title_entry.get().strip()
        content = self.content_text.get("1.0", tk.END).strip()
        slide_type = self.slide_type.get()

        if not title or not content:
            self.status_label.config(text="⚠️ कृपया टाइटल और कंटेंट भरें!", fg="red")
            return

        self.slides.append((title, content, slide_type))  

        if slide_type == "code":
            self.add_code_slide(title, content)
        else:
            self.add_text_slide(title, content)

        # GUI में स्टेटस दिखाएं
        self.status_label.config(text="✅ स्लाइड सफलतापूर्वक जोड़ी गई!", fg="green")

        # 🔥 स्लाइड जोड़ने के बाद इनपुट फील्ड क्लियर करें
        self.title_entry.delete(0, tk.END)
        self.content_text.delete("1.0", tk.END)

        # अपडेट करें कि हम नई स्लाइड पर आ गए हैं
        self.current_slide = len(self.slides) - 1  

        # नेविगेशन बटन को अपडेट करें
        if len(self.slides) > 1:
            self.next_button.config(state=tk.NORMAL)

    def add_text_slide(self, title, content):
        """सिंपल टेक्स्ट स्लाइड ऐड करें"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

    def add_code_slide(self, title, code):
        """कोड स्लाइड ऐड करें"""
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        wrapped_code = textwrap.fill(code, width=80)
        text_frame.text = wrapped_code

    def next_slide(self):
        """अगली स्लाइड पर जाएं"""
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            self.show_slide()

    def previous_slide(self):
        """पिछली स्लाइड पर जाएं"""
        if self.current_slide > 0:
            self.current_slide -= 1
            self.show_slide()

    def show_slide(self):
        """मौजूदा स्लाइड डेटा GUI में दिखाएं"""
        if 0 <= self.current_slide < len(self.slides):
            title, content, slide_type = self.slides[self.current_slide]
            self.title_entry.delete(0, tk.END)
            self.title_entry.insert(0, title)

            self.slide_type.set(slide_type)

            self.content_text.delete("1.0", tk.END)
            self.content_text.insert("1.0", content)

        # नेविगेशन बटन एनेबल/डिसेबल करें
        if self.current_slide == 0:
            self.prev_button.config(state=tk.DISABLED)
        else:
            self.prev_button.config(state=tk.NORMAL)

        if self.current_slide == len(self.slides) - 1:
            self.next_button.config(state=tk.DISABLED)
        else:
            self.next_button.config(state=tk.NORMAL)

    def save_and_exit(self):
        """PPT को सेव करें और प्रोग्राम बंद करें"""
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                 filetypes=[("PowerPoint Files", "*.pptx")],
                                                 title="PPT को सेव करने के लिए स्थान चुनें")
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(text=f"✅ PPT सेव हो गई: {save_path}", fg="blue")
            self.root.quit()

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()
