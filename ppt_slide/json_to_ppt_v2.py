import os
import json
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")
        
        self.prs = Presentation()
        
        tk.Button(root, text="📂 Load JSON & Create PPT", font=("Arial", 12), command=self.load_json).pack(pady=20)
        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.pack()
    
    def load_json(self):
        file_path = filedialog.askopenfilename(filetypes=[("JSON Files", "*.json")], title="Select JSON File")
        if not file_path:
            return
        
        with open(file_path, "r", encoding="utf-8") as file:
            slides_data = json.load(file)
            
        for slide in slides_data:
            title = slide.get("title", "Untitled")
            content = slide.get("content", "")
            slide_type = slide.get("slide_type", "text")
            font_size = slide.get("font_size", 16)
            font_color = slide.get("font_color", "#FFFFFF")  # Default white text
            background_color = slide.get("background_color", "#000000" if slide_type == "code" else "#FFFFFF")
            
            if slide_type == "code":
                self.add_code_slide(title, content, font_size, font_color, background_color)
            else:
                self.add_text_slide(title, content)
        
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")], title="Save PPT")
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(text=f"✅ PPT saved: {save_path}", fg="blue")

    def add_text_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = content

    def add_code_slide(self, title, code, font_size, text_color, bg_color):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        formatted_code = self.format_code_for_ppt(code, "java")
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False

        p = text_frame.add_paragraph()
        p.text = formatted_code
        p.font.name = "Courier New"
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))

    def format_code_for_ppt(self, code, language):
        try:
            lexer = lexers.get_lexer_by_name(language)
        except pygments.util.ClassNotFound:
            lexer = lexers.get_lexer_by_name("python")
        formatter = formatters.HtmlFormatter()
        highlighted_code = pygments.highlight(code, lexer, formatter)
        soup = BeautifulSoup(highlighted_code, "html.parser")
        return soup.get_text()

    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()

# C:\Users\dhruv\Videos\Video_Extractor_Python_Files\python_video_extractor_files\ppt_slide



# To make this PowerPoint content generation more productive and automated, consider the following enhancements:  

# ### **1. Automate Font Size and Styling Adjustments**  
# - Use a dynamic font size based on content length.  
# - Implement auto-resizing if the text exceeds a certain limit.  
# - Adjust font sizes automatically using `pptx.util.Pt()` based on slide size.  

# ### **2. Smarter Code Formatting & Highlighting**  
# - Use **highlight.js** (in combination with `pygments`) to improve syntax highlighting.  
# - Store and apply different themes for better visual appeal.  
# - Render code as images for better preservation of formatting.  

# ### **3. Add AI-Powered Summarization**  
# - Use **GPT-based models** (like OpenAI's API) to summarize lengthy text content.  
# - Automatically shorten long text while keeping it meaningful.  

# ### **4. Optimize Slide Layouts Dynamically**  
# - Choose layouts based on the text length.  
# - Distribute content into multiple slides if needed.  
# - Create title slides, bullet points, and summary slides automatically.  

# ### **5. Automate Background & Theme Selection**  
# - Randomly apply a set of predefined themes.  
# - Detect dark/light mode based on content type.  
# - Allow users to choose styles from a JSON config file.  

# ### **6. Batch Processing & Scheduling**  
# - Allow users to upload multiple JSON files at once.  
# - Schedule PowerPoint generation using a cron job.  
# - Enable bulk processing and auto-save generated PPTs.  

# Would you like me to modify your script to implement some of these?