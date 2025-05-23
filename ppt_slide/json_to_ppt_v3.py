import os
import json
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers, formatters
from bs4 import BeautifulSoup
import textwrap

class PPTCreatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Generator")
        
        self.prs = Presentation()
        
        tk.Button(root, text="📂 Load JSON & Create PPT", font=("Arial", 12), command=self.load_json).pack(pady=20)
        self.status_label = tk.Label(root, text="", font=("Arial", 10), fg="green")
        self.status_label.pack()
    
    def load_json(self):
        file_path = filedialog.askopenfilename(filetypes=[("JSON Files", "*.json")], title="Select JSON File")
        if not file_path:
            return
        
        with open(file_path, "r", encoding="utf-8") as file:
            slides_data = json.load(file)
            
        for slide in slides_data:
            title = slide.get("title", "Untitled")
            content = slide.get("content", "")
            slide_type = slide.get("slide_type", "text")
            font_color = slide.get("font_color", "#FFFFFF")  # Default white text
            background_color = slide.get("background_color", "#000000" if slide_type == "code" else "#FFFFFF")
            
            if slide_type == "code":
                self.add_code_slide(title, content, font_color, background_color)
            else:
                self.add_text_slide(title, content)
        
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")], title="Save PPT")
        if save_path:
            self.prs.save(save_path)
            self.status_label.config(text=f"✅ PPT saved: {save_path}", fg="blue")

    def add_text_slide(self, title, content):
            max_words_per_slide = 300  # Adjust for readability
            words = content.split()

            for i in range(0, len(words), max_words_per_slide):
                slide_layout = self.prs.slide_layouts[1]  # Title + Content layout
                slide = self.prs.slides.add_slide(slide_layout)

                # Ensure title placeholder exists before using it
                if slide.shapes.title:
                    slide.shapes.title.text = title if i == 0 else f"{title} (contd.)"

                # Remove any empty placeholders to avoid "Click to add text" issue
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip() == "":
                        sp = shape
                        slide.shapes._spTree.remove(sp._element)

                # Add a new textbox for content
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
                text_frame = textbox.text_frame
                text_frame.word_wrap = True

                p = text_frame.add_paragraph()
                p.text = " ".join(words[i:i + max_words_per_slide])
                p.font.name = "Mangal"  # Unicode Hindi-compatible font
                p.font.size = Pt(20)  # Adjust size for readability



    def add_code_slide(self, title, code, text_color, bg_color):
                max_lines_per_slide = 20  # Adjust as needed
                lines = code.split("\n")
                
                for i in range(0, len(lines), max_lines_per_slide):
                    slide_layout = self.prs.slide_layouts[5]
                    slide = self.prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = title if i == 0 else f"{title} (contd.)"

                    code_chunk = "\n".join(lines[i:i+max_lines_per_slide])
                    formatted_code = self.format_code_for_ppt(code_chunk, "java")

                    font_size = self.calculate_dynamic_font_size(code_chunk)
                    
                    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(6))
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = False

                    p = text_frame.add_paragraph()
                    p.text = formatted_code
                    p.font.name = "Courier New"
                    p.font.size = Pt(font_size)
                    p.alignment = PP_ALIGN.LEFT
                    p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))


                    textbox.fill.solid()
                    textbox.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))

    def format_code_for_ppt(self, code, language):
            try:
                lexer = lexers.get_lexer_by_name(language)
            except pygments.util.ClassNotFound:
                lexer = lexers.get_lexer_by_name("python")
            formatter = formatters.HtmlFormatter()
            highlighted_code = pygments.highlight(code, lexer, formatter)
            soup = BeautifulSoup(highlighted_code, "html.parser")
            return soup.get_text()

    def calculate_dynamic_font_size(self, code):
            line_count = code.count('\n') + 1
            if line_count < 10:
                return 24
            elif line_count < 20:
                return 18
            elif line_count < 30:
                return 14
            else:
                return 12

    def hex_to_rgb(self, hex_color):            



       
            hex_color = hex_color.lstrip("#")
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

if __name__ == "__main__":
    root = tk.Tk()
    app = PPTCreatorApp(root)
    root.mainloop()
